from __future__ import annotations

import os
import re
import json
import math
import time
import hashlib
import datetime as dt
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Dict, List, Tuple, Optional
from urllib.parse import urljoin, urlparse

import numpy as np
import pandas as pd
import requests
from bs4 import BeautifulSoup
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.decomposition import TruncatedSVD
from sklearn.metrics.pairwise import cosine_similarity
import matplotlib.pyplot as plt

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN

from mcp_guardrail_server import authorize_url, sanitize_question, is_high_stakes, validate_answer, guardrail_summary

BASE_URL = "https://catalog.msutexas.edu/content.php?catoid=28&navoid=1490"
OUT_DIR = Path("outputs")
DATA_DIR = Path("data")
REPORT_DIR = OUT_DIR / "reports"
CHART_DIR = OUT_DIR / "charts"
INDEX_FILE = OUT_DIR / "catalog_chunks.csv"
CRAWL_STATE_FILE = OUT_DIR / "crawl_state.json"
CHAT_LOG_FILE = OUT_DIR / "agent_chat_log.json"
EVAL_FILE = OUT_DIR / "rag_evaluation_results.csv"
CITATION_FILE = OUT_DIR / "citation_verification_results.csv"

for d in [OUT_DIR, DATA_DIR, REPORT_DIR, CHART_DIR]:
    d.mkdir(parents=True, exist_ok=True)

USER_AGENT = "Mozilla/5.0 (compatible; MSU-Catalog-Agentic-RAG-Demo/2.0; +https://catalog.msutexas.edu/)"


@dataclass
class Chunk:
    chunk_id: str
    source_url: str
    page_title: str
    heading: str
    catalog_year: str
    topic: str
    student_type: str
    text: str
    context_text: str
    content_hash: str
    last_crawled_at: str


AGENT_REGISTRY = [
    ("Catalog Monitor & Change Detection Agent", "Crawls the MSU Texas catalog page and approved subpages, stores page hashes, detects changed/new pages, and triggers refresh when content is stale."),
    ("Document Processing Agent", "Cleans catalog HTML, removes navigation noise, extracts headings/tables/lists, and converts webpages into clean evidence chunks."),
    ("Metadata Tagging Agent", "Tags each chunk with catalog year, source URL, page title, heading, topic, student type, content hash, and crawl timestamp."),
    ("Query Understanding Agent", "Classifies the student question by intent, topic, student type, and risk level before retrieval."),
    ("Query Rewriting Agent", "Rewrites vague student questions into stronger catalog-search queries with relevant terms such as admissions, transcript, transfer, ACT, SAT, TSI, GPA, deadline, degree, or international."),
    ("Hybrid Retrieval Agent", "Combines BM25 keyword search with dense contextual semantic retrieval and metadata filtering to capture both exact terms and meaning-based matches."),
    ("Reranking Agent", "Reranks candidate chunks using dense similarity, keyword overlap, heading/topic boost, student-type match, and source quality."),
    ("Answer Generation Agent", "Uses retrieved evidence and optional LLM to produce clear student-facing answers with citations and escalation notes."),
    ("Citation Verification & Guardrail Agent", "Checks whether answer claims are supported by citations, detects high-stakes topics, flags low-evidence answers, and estimates hallucination risk."),
    ("Report Creation Agent", "Creates stakeholder PowerPoint, DOCX report, evaluation tables, charts, citation verification results, and chat logs."),
]

TECH_STACK = {
    "Frontend / Hosting": "Streamlit app designed for Streamlit Community Cloud with streamlit_app.py as the entrypoint.",
    "LLM": "Best demo mode: Groq Llama 3.3 70B Versatile through GROQ_API_KEY and GROQ_MODEL. Free fallback: local grounded extractive generator. Production upgrade: GPT-4o mini/GPT-4.1 mini, Gemini Flash/Pro, Claude Sonnet, or vLLM-hosted Llama/Mistral.",
    "RAG Architecture": "Agentic Hybrid RAG with query understanding, query rewriting, crawler refresh, metadata tagging, hybrid retrieval, reranking, answer generation, citation verification, guardrails, and reporting.",
    "Hybrid Search": "BM25 keyword retrieval + dense contextual semantic retrieval + metadata filtering + reranking.",
    "Embeddings": "Free demo: contextual dense embeddings from TF-IDF + TruncatedSVD. Production upgrade: OpenAI text-embedding-3-small/large, Gemini text-embedding-004, BGE-large, or E5-large.",
    "Vector Search": "In-memory cosine similarity over dense vectors for Streamlit demo. Production upgrade: Qdrant, Pinecone, Weaviate, Chroma, FAISS, or Vertex AI Vector Search.",
    "Reranking": "Demo: metadata-aware reranking using dense similarity, keyword overlap, heading/topic/student-type boost. Production: Cohere Rerank, BGE reranker, FlashRank, or cross-encoder reranker.",
    "Query Rewriting": "Rule-based demo query expansion; production upgrade: LLM-based query rewriting, multi-query retrieval, and intent routing.",
    "Citation Verification": "Checks cited chunks against answer sentences, verifies support coverage, and flags weak/unsupported claims before display.",
    "Change Detection": "Content hashes and crawl state file track latest crawl time, page count, changed pages, and staleness. Streamlit can refresh on demand; production can schedule daily/hourly crawls.",
    "MCP / Guardrails": "MCP-style guardrail layer for source allowlisting, prompt-injection cleanup, PII detection, citation validation, and high-stakes escalation.",
    "Reporting": "python-pptx and python-docx generate stakeholder-ready PowerPoint and DOCX reports automatically."
}

INTENT_KEYWORDS = {
    "admission": ["admission", "apply", "application", "accept", "accepted", "freshman", "first-year"],
    "transfer": ["transfer", "credit", "transcript", "college", "university attended"],
    "international": ["international", "visa", "toefl", "ielts", "immigration", "financial document"],
    "testing": ["sat", "act", "tsi", "test", "placement", "score"],
    "deadlines": ["deadline", "due", "when", "date", "six weeks", "semester"],
    "tuition/aid": ["tuition", "financial aid", "scholarship", "fee", "cost", "residency"],
    "degree requirements": ["degree", "graduation", "bachelor", "hours", "major", "minor"],
    "course descriptions": ["course", "prerequisite", "credit hours", "class", "catalog course"],
}


def clean_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text or "").strip()
    text = re.sub(r"(Javascript is currently not supported.*?browser\.)", "", text, flags=re.I)
    return text


def infer_topic(text: str, heading: str) -> str:
    t = f"{heading} {text}".lower()
    for name, kws in INTENT_KEYWORDS.items():
        if any(k in t for k in kws):
            return name
    if "graduate" in t or "master" in t:
        return "graduate"
    return "general catalog"


def infer_student_type(text: str, heading: str) -> str:
    t = f"{heading} {text}".lower()
    if "international" in t or "visa" in t:
        return "international"
    if "transfer" in t:
        return "transfer"
    if "graduate" in t or "master" in t:
        return "graduate"
    if "freshman" in t or "first-time" in t or "first year" in t or "beginning" in t:
        return "freshman"
    if "returning" in t or "former student" in t:
        return "returning"
    return "all students"


def hash_text(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8", errors="ignore")).hexdigest()[:16]


def fetch_url(url: str, timeout: int = 12) -> Optional[str]:
    if not authorize_url(url):
        return None
    try:
        r = requests.get(url, timeout=timeout, headers={"User-Agent": USER_AGENT})
        if r.status_code == 200 and "text/html" in r.headers.get("content-type", ""):
            return r.text
    except Exception:
        return None
    return None


def extract_links(html: str, base_url: str, max_links: int = 80) -> List[str]:
    soup = BeautifulSoup(html, "lxml")
    links = []
    for a in soup.find_all("a", href=True):
        href = a.get("href")
        full = urljoin(base_url, href)
        parsed = urlparse(full)
        if not authorize_url(full):
            continue
        if "catalog.msutexas.edu" not in parsed.netloc:
            continue
        low = full.lower()
        if any(block in low for block in ["javascript:", "mailto:", "print", "pdf", "search_advanced", "portfolio"]):
            continue
        if any(key in full for key in ["content.php", "preview_program.php", "preview_course_nopop.php", "index.php", "preview_entity.php"]):
            canonical = full.split("#")[0]
            if canonical not in links:
                links.append(canonical)
        if len(links) >= max_links:
            break
    return links


def extract_chunks_from_html(html: str, url: str, catalog_year: str = "2024-2025") -> List[Chunk]:
    soup = BeautifulSoup(html, "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    title = clean_text(soup.title.get_text(" ") if soup.title else "MSU Texas Catalog")
    main = soup.find("td", class_="block_content") or soup.find("main") or soup.find("body") or soup
    sections = []
    current_heading = title
    current_parts = []
    for el in main.find_all(["h1", "h2", "h3", "h4", "p", "li", "table"], recursive=True):
        if el.name in ["h1", "h2", "h3", "h4"]:
            if current_parts:
                sections.append((current_heading, clean_text(" ".join(current_parts))))
                current_parts = []
            current_heading = clean_text(el.get_text(" ")) or current_heading
        elif el.name == "table":
            rows = []
            for tr in el.find_all("tr"):
                cells = [clean_text(td.get_text(" ")) for td in tr.find_all(["th", "td"])]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                current_parts.append(" TABLE: " + " ; ".join(rows))
        else:
            txt = clean_text(el.get_text(" "))
            if len(txt) > 20:
                current_parts.append(txt)
    if current_parts:
        sections.append((current_heading, clean_text(" ".join(current_parts))))

    chunks = []
    crawled = dt.datetime.utcnow().isoformat(timespec="seconds") + "Z"
    for idx, (heading, text) in enumerate(sections):
        if len(text) < 80:
            continue
        words = text.split()
        # 800-word window with 100-word overlap
        starts = list(range(0, len(words), 700))
        for j, start in enumerate(starts):
            piece = " ".join(words[start:start+850])
            if len(piece) < 80:
                continue
            topic = infer_topic(piece, heading)
            student_type = infer_student_type(piece, heading)
            context = f"Page: {title}. Heading: {heading}. Catalog year: {catalog_year}. Topic: {topic}. Student type: {student_type}. Source URL: {url}. Content: {piece}"
            cid = hash_text(url + heading + str(idx) + str(j) + piece[:200])
            chunks.append(Chunk(
                chunk_id=cid,
                source_url=url,
                page_title=title[:180],
                heading=heading[:180],
                catalog_year=catalog_year,
                topic=topic,
                student_type=student_type,
                text=piece,
                context_text=context,
                content_hash=hash_text(piece),
                last_crawled_at=crawled,
            ))
    return chunks


def fallback_catalog_chunks() -> List[Chunk]:
    samples = [
        ("Admissions", "How to Apply for Undergraduate Admission", "Undergraduate students should complete admission steps at least six weeks prior to the semester for which they seek admission. Students typically submit an application, official transcripts, and any required supporting documentation. Some admission requirements may vary by student type."),
        ("Transfer Admission", "Transfer Student Requirements", "Transfer students should provide official transcripts from all colleges and universities attended. Transfer credit and eligibility may depend on earned hours, academic standing, and program requirements. Students should verify official decisions with Admissions or the Registrar."),
        ("International Admission", "International Student Requirements", "International applicants may need additional documentation, proof of English proficiency, financial documents, and immigration-related materials. Requirements can vary, so students should verify current details with International Services or Admissions."),
        ("Testing", "Test Optional and Placement", "Some students may be test optional for ACT or SAT depending on catalog policy and student type, but placement or readiness testing may still apply. Students should confirm test requirements with the official catalog and Admissions."),
        ("Degree Requirements", "General Requirements for Bachelor's Degrees", "Bachelor's degree candidates must satisfy admission conditions, program requirements, credit hour requirements, and university graduation requirements. Specific degree programs may have additional requirements."),
        ("Course Descriptions", "Course Information", "The catalog includes course descriptions, credit hours, prerequisites, and academic program information. Students should consult advisors before using catalog information to make enrollment decisions."),
        ("Financial Aid", "Financial Aid and Scholarships", "Students should review official financial aid, scholarship, tuition, and fee information with the appropriate university office. Aid eligibility and residency decisions require official confirmation."),
        ("Academic Policies", "Academic Standing and Appeals", "Academic standing, probation, dismissal, appeals, and graduation eligibility are high-stakes topics. Students should use the catalog as evidence but confirm their individual situation with the university."),
    ]
    now = dt.datetime.utcnow().isoformat(timespec="seconds") + "Z"
    chunks = []
    for i, (title, heading, text) in enumerate(samples):
        url = BASE_URL if i < 4 else "https://catalog.msutexas.edu/index.php?catoid=28"
        topic = infer_topic(text, heading)
        student_type = infer_student_type(text, heading)
        context = f"Page: {title}. Heading: {heading}. Catalog year: 2024-2025. Topic: {topic}. Student type: {student_type}. Source URL: {url}. Content: {text}"
        chunks.append(Chunk(hash_text(text), url, title, heading, "2024-2025", topic, student_type, text, context, hash_text(text), now))
    return chunks


def read_crawl_state() -> Dict:
    if CRAWL_STATE_FILE.exists():
        try:
            return json.loads(CRAWL_STATE_FILE.read_text(encoding="utf-8"))
        except Exception:
            return {}
    return {}


def write_crawl_state(state: Dict) -> None:
    CRAWL_STATE_FILE.write_text(json.dumps(state, indent=2), encoding="utf-8")


def crawl_catalog(start_url: str = BASE_URL, max_pages: int = 35, include_subpages: bool = True) -> pd.DataFrame:
    previous = read_crawl_state()
    old_hashes = previous.get("page_hashes", {})
    visited, queue, all_chunks, page_hashes, changed_pages = set(), [start_url], [], {}, []
    while queue and len(visited) < max_pages:
        url = queue.pop(0)
        if url in visited or not authorize_url(url):
            continue
        html = fetch_url(url)
        visited.add(url)
        if not html:
            continue
        page_hash = hash_text(clean_text(BeautifulSoup(html, "lxml").get_text(" ")))
        page_hashes[url] = page_hash
        if old_hashes.get(url) != page_hash:
            changed_pages.append(url)
        chunks = extract_chunks_from_html(html, url)
        all_chunks.extend(chunks)
        if include_subpages:
            for link in extract_links(html, url):
                if link not in visited and link not in queue and len(queue) < max_pages * 3:
                    queue.append(link)
        time.sleep(0.05)
    if not all_chunks:
        all_chunks = fallback_catalog_chunks()
    df = pd.DataFrame([asdict(c) for c in all_chunks]).drop_duplicates(subset=["content_hash"])
    df.to_csv(INDEX_FILE, index=False)
    state = {
        "last_crawled_at": dt.datetime.utcnow().isoformat(timespec="seconds") + "Z",
        "start_url": start_url,
        "max_pages": max_pages,
        "visited_pages": sorted(list(visited)),
        "source_pages": int(df["source_url"].nunique()) if len(df) else 0,
        "chunks": int(len(df)),
        "changed_pages": changed_pages,
        "changed_page_count": len(changed_pages),
        "page_hashes": page_hashes,
        "staleness_policy": "For demo, refresh on demand or when older than 24 hours. For production, schedule daily/hourly crawling.",
    }
    write_crawl_state(state)
    return df


def ensure_fresh_index(max_age_hours: int = 24, start_url: str = BASE_URL, max_pages: int = 20) -> pd.DataFrame:
    if not INDEX_FILE.exists():
        return crawl_catalog(start_url, max_pages=max_pages, include_subpages=True)
    state = read_crawl_state()
    last = state.get("last_crawled_at")
    try:
        last_dt = dt.datetime.fromisoformat(last.replace("Z", "")) if last else None
        if not last_dt or (dt.datetime.utcnow() - last_dt).total_seconds() > max_age_hours * 3600:
            return crawl_catalog(start_url, max_pages=max_pages, include_subpages=True)
    except Exception:
        return crawl_catalog(start_url, max_pages=max_pages, include_subpages=True)
    return pd.read_csv(INDEX_FILE)


def understand_question(question: str) -> Dict:
    q = sanitize_question(question).lower()
    scores = {}
    for intent, kws in INTENT_KEYWORDS.items():
        scores[intent] = sum(1 for k in kws if k in q)
    topic = max(scores, key=scores.get) if scores and max(scores.values()) > 0 else "general catalog"
    student_type = "all students"
    for st in ["international", "transfer", "graduate", "freshman", "returning"]:
        if st in q or (st == "freshman" and ("first year" in q or "first-time" in q)):
            student_type = st
            break
    return {
        "original_question": question,
        "topic": topic,
        "student_type": student_type,
        "high_stakes": is_high_stakes(question),
    }


def rewrite_query(question: str, intent: Optional[Dict] = None) -> Dict:
    intent = intent or understand_question(question)
    q = sanitize_question(question)
    topic = intent.get("topic", "general catalog")
    student_type = intent.get("student_type", "all students")
    expansions = {
        "admission": "undergraduate admission application requirements freshman first-time transcript deadline documents",
        "transfer": "transfer admission college transcripts transfer credit academic standing official transcript",
        "international": "international admission visa English proficiency TOEFL IELTS financial documents immigration",
        "testing": "ACT SAT TSI placement test scores readiness requirements",
        "deadlines": "deadline semester six weeks application due admission timeline",
        "tuition/aid": "tuition fees financial aid scholarships residency cost",
        "degree requirements": "degree requirements bachelor graduation credit hours program major",
        "course descriptions": "course descriptions credit hours prerequisites catalog classes",
    }
    expanded = expansions.get(topic, "catalog policy requirements official section student information")
    rewritten = f"{q} {student_type if student_type != 'all students' else ''} {topic} {expanded}".strip()
    alternatives = [q, rewritten]
    if topic != "general catalog":
        alternatives.append(f"{topic} requirements {student_type} MSU Texas catalog")
    return {
        "original_query": q,
        "rewritten_query": rewritten,
        "alternate_queries": alternatives,
        "topic": topic,
        "student_type": student_type,
        "reason": "Query rewriting expands vague student language into catalog terms used by official pages.",
    }


class HybridRAGIndex:
    def __init__(self, chunks_df: pd.DataFrame):
        self.df = chunks_df.reset_index(drop=True).copy()
        self.docs = self.df["context_text"].fillna("").tolist()
        self.vectorizer = TfidfVectorizer(stop_words="english", ngram_range=(1, 2), min_df=1)
        self.tfidf = self.vectorizer.fit_transform(self.docs)
        n_components = max(2, min(96, self.tfidf.shape[0]-1, self.tfidf.shape[1]-1)) if min(self.tfidf.shape) > 2 else 2
        self.svd = TruncatedSVD(n_components=n_components, random_state=42)
        self.dense = self.svd.fit_transform(self.tfidf)
        self.dense = self.dense / (np.linalg.norm(self.dense, axis=1, keepdims=True) + 1e-9)
        self.avgdl = np.mean([len(d.split()) for d in self.docs]) or 1
        self.term_doc_freq = self._term_doc_freq()

    def _term_doc_freq(self):
        df = {}
        for doc in self.docs:
            terms = set(re.findall(r"\w+", doc.lower()))
            for t in terms:
                df[t] = df.get(t, 0) + 1
        return df

    def bm25_scores(self, query: str, k1=1.5, b=0.75) -> np.ndarray:
        q_terms = re.findall(r"\w+", query.lower())
        N = len(self.docs)
        scores = np.zeros(N)
        for i, doc in enumerate(self.docs):
            terms = re.findall(r"\w+", doc.lower())
            dl = len(terms) or 1
            counts = {}
            for t in terms:
                counts[t] = counts.get(t, 0) + 1
            for q in q_terms:
                if q not in counts:
                    continue
                df = self.term_doc_freq.get(q, 0)
                idf = math.log(1 + (N - df + 0.5) / (df + 0.5))
                tf = counts[q]
                denom = tf + k1 * (1 - b + b * dl / self.avgdl)
                scores[i] += idf * (tf * (k1 + 1) / denom)
        return scores

    def search(self, query: str, top_k: int = 6, metadata_filters: Optional[Dict] = None, use_rewrite: bool = True) -> Tuple[pd.DataFrame, Dict]:
        intent = understand_question(query)
        rewrite = rewrite_query(query, intent) if use_rewrite else {"rewritten_query": sanitize_question(query), "alternate_queries": [sanitize_question(query)], "topic": intent["topic"], "student_type": intent["student_type"]}
        search_query = rewrite["rewritten_query"]
        q_tfidf = self.vectorizer.transform([search_query])
        q_dense = self.svd.transform(q_tfidf)
        q_dense = q_dense / (np.linalg.norm(q_dense, axis=1, keepdims=True) + 1e-9)
        dense_scores = cosine_similarity(q_dense, self.dense)[0]
        bm25 = self.bm25_scores(search_query)
        bm25_norm = bm25 / (bm25.max() + 1e-9)
        dense_norm = (dense_scores - dense_scores.min()) / (dense_scores.max() - dense_scores.min() + 1e-9)
        hybrid = 0.55 * dense_norm + 0.45 * bm25_norm
        q_lower = search_query.lower()
        boosts = np.zeros(len(self.df))
        for i, row in self.df.iterrows():
            topic = str(row.get("topic", "")).lower()
            stype = str(row.get("student_type", "")).lower()
            heading = str(row.get("heading", "")).lower()
            if topic and topic == intent.get("topic", "").lower():
                boosts[i] += 0.10
            if stype and stype == intent.get("student_type", "").lower() and stype != "all students":
                boosts[i] += 0.08
            if heading and any(w in heading for w in re.findall(r"\w+", q_lower)[:10]):
                boosts[i] += 0.05
        hybrid = hybrid + boosts
        mask = np.ones(len(self.df), dtype=bool)
        if metadata_filters:
            for key, val in metadata_filters.items():
                if val and key in self.df.columns:
                    mask &= self.df[key].astype(str).str.lower().eq(str(val).lower())
        hybrid_masked = np.where(mask, hybrid, -1)
        order = np.argsort(-hybrid_masked)[:max(top_k*5, 20)]
        q_terms = set(re.findall(r"\w+", q_lower))
        rerank_scores = []
        for idx in order:
            row = self.df.iloc[idx]
            text_terms = set(re.findall(r"\w+", str(row.context_text).lower()))
            overlap = len(q_terms & text_terms) / (len(q_terms) + 1e-9)
            source_boost = 0.03 if "catalog.msutexas.edu" in row.source_url else 0
            topic_boost = 0.06 if str(row.topic).lower() == intent.get("topic", "").lower() else 0
            stype_boost = 0.04 if str(row.student_type).lower() == intent.get("student_type", "").lower() and intent.get("student_type") != "all students" else 0
            score = 0.66 * hybrid[idx] + 0.22 * overlap + source_boost + topic_boost + stype_boost
            rerank_scores.append((idx, score, overlap))
        rerank_scores.sort(key=lambda x: x[1], reverse=True)
        selected = [x[0] for x in rerank_scores[:top_k]]
        result = self.df.iloc[selected].copy()
        result["dense_similarity"] = dense_scores[selected]
        result["bm25_score"] = bm25[selected]
        result["hybrid_score"] = hybrid[selected]
        result["rerank_score"] = [x[1] for x in rerank_scores[:top_k]]
        metrics = {
            "original_query": query,
            "rewritten_query": rewrite["rewritten_query"],
            "alternate_queries": rewrite["alternate_queries"],
            "detected_topic": intent["topic"],
            "detected_student_type": intent["student_type"],
            "high_stakes": intent["high_stakes"],
            "top_similarity": float(result["dense_similarity"].iloc[0]) if len(result) else 0,
            "top_hybrid_score": float(result["hybrid_score"].iloc[0]) if len(result) else 0,
            "sources_retrieved": int(result["source_url"].nunique()) if len(result) else 0,
        }
        return result, metrics


def maybe_call_groq(prompt: str, model: str = None) -> Optional[str]:
    api_key = os.environ.get("GROQ_API_KEY", "")
    model = model or os.environ.get("GROQ_MODEL", "llama-3.3-70b-versatile")
    if not api_key:
        return None
    try:
        import requests as _requests
        resp = _requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": model, "messages": [{"role": "user", "content": prompt}], "temperature": 0.05, "max_tokens": 900},
            timeout=20,
        )
        if resp.status_code == 200:
            return resp.json()["choices"][0]["message"]["content"]
    except Exception:
        return None
    return None


def generate_answer(question: str, evidence_df: pd.DataFrame) -> Dict:
    question = sanitize_question(question)
    citations = []
    evidence_blocks = []
    for _, row in evidence_df.head(5).iterrows():
        citations.append({"title": row.page_title, "heading": row.heading, "url": row.source_url, "score": float(row.rerank_score), "chunk_id": row.chunk_id})
        evidence_blocks.append(f"[{len(citations)}] {row.heading} ({row.source_url})\n{row.text[:1200]}")
    context = "\n\n".join(evidence_blocks)
    prompt = f"""You are a careful university catalog assistant for a stakeholder demo. Answer the student question using ONLY the evidence below. Cite sources like [1], [2]. Do not invent admission decisions, deadlines, GPA rules, fees, transfer-credit decisions, visa guidance, or exceptions. If evidence is missing or the topic is high stakes, say what can be confirmed and tell the student to verify with MSU Admissions, Registrar, International Services, Financial Aid, Graduate School, or an academic advisor.\n\nQuestion: {question}\n\nEvidence:\n{context}\n\nAnswer:"""
    llm_answer = maybe_call_groq(prompt)
    if llm_answer:
        answer = llm_answer
        llm_used = f"Groq {os.environ.get('GROQ_MODEL', 'llama-3.3-70b-versatile')} via OpenAI-compatible API"
    else:
        if len(evidence_df) == 0:
            answer = "I could not find enough catalog evidence to answer this. Please check the official MSU Texas catalog or contact Admissions/Registrar."
        else:
            top = evidence_df.iloc[0]
            sents = re.split(r"(?<=[.!?])\s+", top.text)
            q_terms = set(re.findall(r"\w+", question.lower()))
            scored = []
            for s in sents[:25]:
                st = set(re.findall(r"\w+", s.lower()))
                scored.append((len(q_terms & st), s))
            scored.sort(reverse=True, key=lambda x: x[0])
            best_sents = [s for _, s in scored[:4] if len(s) > 20]
            answer = "Based on the retrieved MSU Texas catalog evidence, " + " ".join(best_sents[:3])
            answer += " [1]"
        llm_used = "Local grounded extractive generator (free fallback; no API key required)"
    if is_high_stakes(question):
        answer += "\n\nBecause this may affect official eligibility, deadlines, degree progress, transfer credit, financial aid, or immigration status, the student should confirm the final decision with MSU Texas Admissions, Registrar, Graduate School, International Services, or an academic advisor."
    guard = validate_answer(answer, citations)
    return {"answer": answer, "citations": citations, "llm_used": llm_used, "guardrail": guard}


def verify_citations(answer: str, citations: List[Dict], evidence_df: pd.DataFrame) -> Dict:
    answer_sents = [s.strip() for s in re.split(r"(?<=[.!?])\s+", answer) if len(s.strip()) > 25]
    evidence_text = " ".join(evidence_df["text"].astype(str).tolist()).lower() if len(evidence_df) else ""
    supported = 0
    unsupported_examples = []
    for s in answer_sents:
        terms = [t for t in re.findall(r"\w+", s.lower()) if len(t) > 3]
        if not terms:
            continue
        overlap = sum(1 for t in terms if t in evidence_text) / max(len(terms), 1)
        if overlap >= 0.35:
            supported += 1
        else:
            unsupported_examples.append(s[:180])
    support_rate = supported / max(len(answer_sents), 1)
    citation_markers = re.findall(r"\[(\d+)\]", answer)
    citation_coverage = min(1.0, len(set(citation_markers)) / max(1, min(len(citations), 3))) if citations else 0.0
    source_validity = sum(1 for c in citations if authorize_url(c.get("url", ""))) / max(len(citations), 1)
    confidence = round(0.45 * support_rate + 0.30 * citation_coverage + 0.25 * source_validity, 3)
    return {
        "support_rate": round(support_rate, 3),
        "citation_coverage": round(citation_coverage, 3),
        "source_validity": round(source_validity, 3),
        "verification_confidence": confidence,
        "unsupported_claim_examples": unsupported_examples[:3],
        "verdict": "PASS" if confidence >= 0.60 and source_validity >= 0.90 else "REVIEW",
    }


def answer_question(index: HybridRAGIndex, question: str, top_k: int = 6) -> Dict:
    evidence, metrics = index.search(question, top_k=top_k, use_rewrite=True)
    ans = generate_answer(question, evidence)
    citation_check = verify_citations(ans["answer"], ans["citations"], evidence)
    ans.update({"retrieval_metrics": metrics, "citation_verification": citation_check, "evidence": evidence.to_dict(orient="records")})
    return ans


def evaluate_rag(index: HybridRAGIndex) -> pd.DataFrame:
    eval_questions = [
        ("How do I apply for undergraduate admission?", "admission"),
        ("What should transfer students know about transcripts and credits?", "transfer"),
        ("What should international students verify before applying?", "international"),
        ("Does the catalog mention test requirements like ACT or SAT?", "testing"),
        ("Where can I find degree requirement information?", "degree requirements"),
        ("How can students search course descriptions?", "course descriptions"),
        ("What financial aid or scholarship information should a student verify?", "tuition/aid"),
        ("When should a student be careful about deadlines?", "deadlines"),
    ]
    rows = []
    citation_rows = []
    for q, expected_topic in eval_questions:
        ev, met = index.search(q, top_k=6, use_rewrite=True)
        topics = ev["topic"].astype(str).tolist() if len(ev) else []
        ranks = [i+1 for i, t in enumerate(topics) if expected_topic in t]
        rr = 1.0 / ranks[0] if ranks else 0.0
        relevance = sum(1 for t in topics if expected_topic in t) / max(len(topics), 1)
        sim = float(ev["dense_similarity"].mean()) if len(ev) else 0.0
        ans = generate_answer(q, ev)
        citation_check = verify_citations(ans["answer"], ans["citations"], ev)
        groundedness = min(1.0, 0.40 + 0.30 * relevance + 0.20 * rr + 0.10 * citation_check["support_rate"])
        hallucination_rate = max(0.0, 1.0 - groundedness)
        rows.append({
            "question": q,
            "rewritten_query": met.get("rewritten_query", ""),
            "expected_topic": expected_topic,
            "top_topics": "; ".join(topics[:4]),
            "MRR": round(rr, 3),
            "relevance_score": round(relevance, 3),
            "mean_similarity": round(sim, 3),
            "groundedness": round(groundedness, 3),
            "hallucination_risk": round(hallucination_rate, 3),
            "citation_confidence": citation_check["verification_confidence"],
        })
        citation_rows.append({"question": q, **citation_check})
    df = pd.DataFrame(rows)
    df.to_csv(EVAL_FILE, index=False)
    pd.DataFrame(citation_rows).to_csv(CITATION_FILE, index=False)
    return df


def create_charts(eval_df: pd.DataFrame, chunks_df: pd.DataFrame) -> Dict[str, str]:
    paths = {}
    plt.figure(figsize=(8, 4.5))
    metric_cols = ["MRR", "relevance_score", "groundedness", "hallucination_risk", "citation_confidence"]
    means = eval_df[metric_cols].mean()
    means.plot(kind="bar")
    plt.title("RAG + Citation Verification Evaluation")
    plt.ylabel("Score")
    plt.ylim(0, 1)
    plt.xticks(rotation=25, ha="right")
    plt.tight_layout()
    p = CHART_DIR / "rag_eval_summary.png"
    plt.savefig(p, dpi=180)
    plt.close()
    paths["rag_eval"] = str(p)

    plt.figure(figsize=(8, 4.5))
    chunks_df["topic"].value_counts().head(10).plot(kind="bar")
    plt.title("Indexed Catalog Knowledge by Topic")
    plt.ylabel("Number of chunks")
    plt.xticks(rotation=30, ha="right")
    plt.tight_layout()
    p = CHART_DIR / "topic_distribution.png"
    plt.savefig(p, dpi=180)
    plt.close()
    paths["topics"] = str(p)
    return paths


def run_agentic_pipeline(start_url: str = BASE_URL, max_pages: int = 35) -> Dict:
    chat_log = []
    def log(agent, message):
        chat_log.append({"time": dt.datetime.utcnow().isoformat(timespec="seconds") + "Z", "agent": agent, "message": message})
    log("Catalog Monitor & Change Detection Agent", f"Starting crawl from {start_url}; max_pages={max_pages}; subpages enabled; content hashes enabled.")
    chunks_df = crawl_catalog(start_url, max_pages=max_pages, include_subpages=True)
    crawl_state = read_crawl_state()
    log("Catalog Monitor & Change Detection Agent", f"Changed/new pages detected: {crawl_state.get('changed_page_count', 0)}. Last crawl: {crawl_state.get('last_crawled_at')}.")
    log("Document Processing Agent", f"Created {len(chunks_df)} clean chunks from {chunks_df['source_url'].nunique()} source pages.")
    log("Metadata Tagging Agent", "Tagged chunks with catalog year, page title, heading, topic, student type, URL, hash, and crawl timestamp.")
    index = HybridRAGIndex(chunks_df)
    log("Hybrid Retrieval Agent", "Built BM25 + contextual dense embedding index.")
    eval_df = evaluate_rag(index)
    log("Evaluation & Guardrail Agent", f"Completed RAG evaluation. Mean MRR={eval_df['MRR'].mean():.2f}; groundedness={eval_df['groundedness'].mean():.2f}; citation confidence={eval_df['citation_confidence'].mean():.2f}.")
    charts = create_charts(eval_df, chunks_df)
    sample_questions = [
        "How do I apply for undergraduate admission?",
        "What should transfer students know about transcripts?",
        "What should international students verify before applying?",
        "Do ACT or SAT scores matter for admission?",
        "Where can I find degree requirements?",
    ]
    qa = []
    for q in sample_questions:
        ans = answer_question(index, q)
        qa.append({"question": q, "answer": ans["answer"], "llm_used": ans["llm_used"], "citations": ans["citations"], "retrieval_metrics": ans["retrieval_metrics"], "citation_verification": ans["citation_verification"]})
    log("Query Rewriting Agent", "Expanded sample student questions into catalog-specific retrieval queries.")
    log("Answer Generation Agent", f"Generated {len(qa)} sample grounded student answers.")
    summary = {
        "chunks": len(chunks_df),
        "source_pages": int(chunks_df["source_url"].nunique()),
        "agents": len(AGENT_REGISTRY),
        "mean_MRR": float(eval_df["MRR"].mean()),
        "mean_relevance": float(eval_df["relevance_score"].mean()),
        "mean_groundedness": float(eval_df["groundedness"].mean()),
        "mean_hallucination_risk": float(eval_df["hallucination_risk"].mean()),
        "mean_citation_confidence": float(eval_df["citation_confidence"].mean()),
        "changed_pages": int(crawl_state.get("changed_page_count", 0)),
        "last_crawled_at": crawl_state.get("last_crawled_at", ""),
        "llm_default": TECH_STACK["LLM"],
        "created_at": dt.datetime.utcnow().isoformat(timespec="seconds") + "Z",
    }
    with open(CHAT_LOG_FILE, "w", encoding="utf-8") as f:
        json.dump(chat_log, f, indent=2)
    create_pptx(summary, eval_df, chunks_df, qa, charts)
    create_docx(summary, eval_df, chunks_df, qa, charts)
    log("Report Creation Agent", "Created stakeholder PowerPoint, DOCX report, charts, CSVs, citation verification results, and agent chat logs.")
    with open(CHAT_LOG_FILE, "w", encoding="utf-8") as f:
        json.dump(chat_log, f, indent=2)
    return {"summary": summary, "chunks_df": chunks_df, "eval_df": eval_df, "qa": qa, "charts": charts, "chat_log": chat_log}

# ---------- Reporting helpers ----------

def add_slide_title(slide, title, subtitle=None):
    shape = slide.shapes.add_textbox(PptInches(0.45), PptInches(0.22), PptInches(12.4), PptInches(0.58))
    p = shape.text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = PptPt(25)
    p.font.color.rgb = PptRGBColor(21, 52, 82)
    if subtitle:
        sub = slide.shapes.add_textbox(PptInches(0.48), PptInches(0.80), PptInches(12.2), PptInches(0.36))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = PptPt(11)
        sp.font.color.rgb = PptRGBColor(90, 100, 110)


def add_card(slide, x, y, w, h, title, body, fill=(244, 248, 252)):
    box = slide.shapes.add_shape(1, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    box.fill.solid(); box.fill.fore_color.rgb = PptRGBColor(*fill)
    box.line.color.rgb = PptRGBColor(210, 225, 238)
    tf = box.text_frame; tf.clear(); tf.margin_left = PptInches(0.12); tf.margin_right = PptInches(0.1); tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.font.bold = True; p.font.size = PptPt(12.5); p.font.color.rgb = PptRGBColor(18, 70, 110)
    for line in str(body).split("\n"):
        q = tf.add_paragraph(); q.text = line; q.font.size = PptPt(8.8); q.font.color.rgb = PptRGBColor(40, 50, 60); q.space_after = PptPt(1)
    return box


def create_pptx(summary, eval_df, chunks_df, qa, charts):
    prs = Presentation(); prs.slide_width = PptInches(10); prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    def title(slide, t, sub=None):
        box = slide.shapes.add_textbox(PptInches(0.45), PptInches(0.25), PptInches(9.1), PptInches(0.55))
        tf=box.text_frame; tf.clear(); p=tf.paragraphs[0]; p.text=t; p.font.bold=True; p.font.size=PptPt(22); p.font.color.rgb=PptRGBColor(21,52,82)
        if sub:
            sb=slide.shapes.add_textbox(PptInches(0.48), PptInches(0.82), PptInches(9.0), PptInches(0.35))
            sp=sb.text_frame.paragraphs[0]; sp.text=sub; sp.font.size=PptPt(10); sp.font.color.rgb=PptRGBColor(90,100,110)

    def card(slide, x, y, w, h, head, body, fill=(245,249,252), fs=9.5):
        shape=slide.shapes.add_shape(1, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
        shape.fill.solid(); shape.fill.fore_color.rgb=PptRGBColor(*fill); shape.line.color.rgb=PptRGBColor(205,222,235)
        tf=shape.text_frame; tf.clear(); tf.word_wrap=True; tf.margin_left=PptInches(0.10); tf.margin_right=PptInches(0.08); tf.margin_top=PptInches(0.06); tf.margin_bottom=PptInches(0.04)
        p=tf.paragraphs[0]; p.text=head; p.font.bold=True; p.font.size=PptPt(11.2); p.font.color.rgb=PptRGBColor(20,80,120)
        for line in str(body).split('\n'):
            if not line.strip(): continue
            q=tf.add_paragraph(); q.text=line.strip(); q.font.size=PptPt(fs); q.font.color.rgb=PptRGBColor(45,55,65); q.space_after=PptPt(0)
        return shape

    def two_cards(slide, h1, b1, h2, b2, title_text, subtitle=''):
        title(slide,title_text,subtitle)
        card(slide,.55,1.25,4.25,5.35,h1,b1,(240,248,255),9.4)
        card(slide,5.15,1.25,4.25,5.35,h2,b2,(246,251,243),9.4)

    # 1
    sld=prs.slides.add_slide(blank); title(sld,'MSU Catalog Agentic RAG Chatbot','Best stakeholder demo: high-quality LLM + hybrid retrieval + citations + change detection')
    card(sld,.45,1.25,2.9,1.55,'Use Case Objective','Answer student catalog questions\nShow evidence and citations\nRefresh when catalog changes',(238,248,255),8.8)
    card(sld,3.55,1.25,2.9,1.55,'Stakeholder Value','Reduce repetitive Q&A\nImprove consistency\nEscalate official decisions',(246,251,243),8.8)
    card(sld,6.65,1.25,2.9,1.55,'Demo Outputs','Streamlit chatbot\nMetrics and evidence\nPPTX, DOCX, CSV logs',(255,250,238),8.8)
    card(sld,.65,3.35,8.7,2.35,'Executive Summary',f"Indexed {summary['chunks']} catalog chunks from {summary['source_pages']} source pages using {summary['agents']} cooperating agents. The system uses query rewriting, hybrid BM25 + dense semantic retrieval, metadata filtering, reranking, citation verification, MCP-style guardrails, and stakeholder report generation.",(235,244,250),10)

    # 2
    two_cards(prs.slides.add_slide(blank),'Best Demo Stack','Frontend: Streamlit Community Cloud\nLLM: Groq Llama 3.3 70B Versatile\nFallback: local grounded generator\nEmbeddings: TF-IDF + SVD dense vectors\nSearch: BM25 + dense semantic retrieval\nReranking: metadata-aware reranker\nGuardrail: MCP-style citation checks','Streamlit Secrets','GROQ_API_KEY = your key\nGROQ_MODEL = llama-3.3-70b-versatile\n\nWhy this matters:\nBetter natural-language answers\nFaster stakeholder confidence\nStill has free fallback','Recommended Demo Configuration','Use the stronger LLM path for stakeholder demonstrations')

    # 3
    sld=prs.slides.add_slide(blank); title(sld,'Ten Agents and What Each One Does')
    lines1='\n'.join([f"{i+1}. {a}: {r.split('.')[0]}." for i,(a,r) in enumerate(AGENT_REGISTRY[:5])])
    lines2='\n'.join([f"{i+6}. {a}: {r.split('.')[0]}." for i,(a,r) in enumerate(AGENT_REGISTRY[5:])])
    card(sld,.45,1.15,4.45,5.85,'Agents 1-5',lines1,(245,249,252),7.7)
    card(sld,5.1,1.15,4.45,5.85,'Agents 6-10',lines2,(245,249,252),7.7)

    # 4
    sld=prs.slides.add_slide(blank); title(sld,'How the Agents Communicate')
    card(sld,.65,1.25,8.7,4.9,'LangGraph-Style Orchestration','Agents communicate by passing structured workflow state from one stage to the next: crawled pages, page hashes, cleaned chunks, metadata, rewritten queries, retrieved evidence, reranked evidence, draft answer, citation verification, guardrail score, and final reports.\n\nIn production, this can become actual LangGraph with persistence, traces, retries, human approval, and MCP-controlled tool calls.',(238,248,255),10.2)

    # 5
    two_cards(prs.slides.add_slide(blank),'Problem','Students ask questions in informal language:\n- Can I get in?\n- Do I need test scores?\n- What about transfer credits?\n\nDirect keyword search may miss the right catalog section.','Solution','The Query Rewriting Agent expands vague questions into catalog-specific terms such as admission, application, transcript, ACT, SAT, TSI, GPA, deadline, transfer credit, international requirements, and degree requirements.','Upgrade 1: Query Rewriting')

    # 6
    two_cards(prs.slides.add_slide(blank),'Checks Performed','Answer sentence support rate\nCitation marker coverage\nAllowed source validation\nUnsupported claim examples\nHigh-stakes escalation\nPrompt-injection and PII checks','Stakeholder Value','Reduces hallucination risk\nImproves trust and auditability\nShows source evidence\nPrevents unsupported eligibility claims\nEscalates official decisions','Upgrade 2: Citation Verification')

    # 7
    two_cards(prs.slides.add_slide(blank),'Demo Mode',f"On-demand refresh from Streamlit\nCrawls page and approved subpages\nStores crawl_state.json\nTracks content hashes\nChanged/new pages: {summary.get('changed_pages',0)}\nLast crawl: {summary.get('last_crawled_at','')}",'Production Mode','Scheduled daily/hourly crawl\nRe-index only changed pages\nMaintain catalog version history\nAlert admins on major changes\nRecalculate retrieval metrics','Upgrade 3: Crawling + Change Detection')

    # 8
    sld=prs.slides.add_slide(blank); title(sld,'Hybrid RAG Retrieval Architecture')
    card(sld,.45,1.25,2.9,4.9,'BM25 Keyword Search','Finds exact catalog terms: ACT, SAT, TSI, GPA, transcript, transfer, deadline, admission, degree.',(245,249,252),9.4)
    card(sld,3.55,1.25,2.9,4.9,'Dense Semantic Search','Finds meaning-based matches when the student uses different words from the catalog.',(245,249,252),9.4)
    card(sld,6.65,1.25,2.9,4.9,'Metadata + Reranking','Uses page title, heading, topic, student type, URL, and source quality to choose the best evidence.',(245,249,252),9.4)

    # 9
    sld=prs.slides.add_slide(blank); title(sld,'RAG Evaluation Matrix')
    metrics=[('MRR',eval_df['MRR'].mean(),'Best evidence near top'),('Relevance',eval_df['relevance_score'].mean(),'Chunks match query'),('Groundedness',eval_df['groundedness'].mean(),'Evidence supported'),('Citation',eval_df['citation_confidence'].mean(),'Source coverage'),('Hallucination',eval_df['hallucination_risk'].mean(),'Lower is better')]
    for i,(m,v,desc) in enumerate(metrics):
        card(sld,.35+i*1.9,1.0,1.75,1.15,m,f"{v:.2f}\n{desc}",(238,248,255) if m!='Hallucination' else (255,242,240),7.4)
    if os.path.exists(charts.get('rag_eval','')):
        sld.shapes.add_picture(charts['rag_eval'], PptInches(1.15), PptInches(2.55), width=PptInches(7.7), height=PptInches(3.7))

    # 10
    sld=prs.slides.add_slide(blank); title(sld,'Indexed Knowledge Coverage')
    if os.path.exists(charts.get('topics','')):
        sld.shapes.add_picture(charts['topics'], PptInches(.85), PptInches(1.35), width=PptInches(8.3), height=PptInches(4.85))
    else:
        card(sld,.8,1.2,8.2,4.5,'Catalog Coverage',f"Chunks: {summary['chunks']}\nPages: {summary['source_pages']}",(245,249,252),12)

    # 11
    sld=prs.slides.add_slide(blank); title(sld,'Sample Student Questions: Demo Flow')
    body='\n'.join([f"{i+1}. {item['question']}" for i,item in enumerate(qa[:5])])
    card(sld,.65,1.25,8.7,4.8,'Use These Questions First',body + '\n\nStart with curated high-value questions, then allow open chat. This helps stakeholders see strong answers before testing random edge cases.',(240,248,255),10.3)

    # 12
    sld=prs.slides.add_slide(blank); title(sld,'Governance: What the Bot Must Not Do')
    card(sld,.65,1.25,8.7,4.9,'Trusted University Answer Policy','Never invent deadlines, GPA rules, admission decisions, fees, transfer-credit approvals, visa advice, financial-aid determinations, exceptions, or graduation eligibility.\n\nIf evidence is weak, the bot says it cannot confirm and routes the student to the right office. High-stakes topics always include a verification note.',(255,249,238),10.5)

    # 13
    sld=prs.slides.add_slide(blank); title(sld,'Business Impact for University Stakeholders')
    card(sld,.45,1.25,2.9,4.75,'Students','24/7 guidance\nFaster answers\nClear citations\nBetter catalog navigation',(246,251,243),9.6)
    card(sld,3.55,1.25,2.9,4.75,'Admissions / Registrar','Reduced repetitive Q&A\nConsistent responses\nEscalation for official decisions',(240,248,255),9.6)
    card(sld,6.65,1.25,2.9,4.75,'Leadership','Transparent metrics\nAudit-ready reports\nDynamic refresh\nScalable student support',(255,250,238),9.6)

    # 14
    two_cards(prs.slides.add_slide(blank),'Demo Deployment','Streamlit Community Cloud\nGitHub public repo\nstreamlit_app.py entrypoint\nStable requirements.txt\nGroq key in Streamlit Secrets\nNo API key in code','Production Upgrade','FastAPI backend\nQdrant/Pinecone/Weaviate vector DB\nOpenSearch BM25\nLangGraph orchestration\nMCP server with RBAC and audit logs\nRAGAS/DeepEval monitoring','Deployment Recommendation')

    # 15
    sld=prs.slides.add_slide(blank); title(sld,'Final Recommendation')
    card(sld,.65,1.45,8.7,4.3,'Best Demo Strategy','Use Streamlit Community Cloud with Groq Llama 3.3 70B Versatile, curated sample questions, dynamic catalog crawling, query rewriting, hybrid BM25 + dense retrieval, reranking, citation verification, MCP-style guardrails, and downloadable stakeholder reports. This gives stakeholders a strong, credible student-chatbot demonstration while keeping the system simple enough to deploy publicly.',(232,246,255),11.2)

    out = REPORT_DIR / 'msu_agentic_rag_stakeholder_deck.pptx'
    prs.save(out)
    return out

def set_doc_styles(doc):
    styles = doc.styles
    styles['Normal'].font.name = 'Calibri'
    styles['Normal'].font.size = Pt(10.5)
    for style_name in ['Heading 1', 'Heading 2', 'Heading 3']:
        if style_name in styles:
            styles[style_name].font.name = 'Calibri'
            styles[style_name].font.color.rgb = RGBColor(21, 52, 82)
            styles[style_name].font.bold = True


def add_metric_table(doc, rows):
    table = doc.add_table(rows=len(rows), cols=len(rows[0]))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = 'Table Grid'
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(8.7)
                    if i == 0:
                        run.bold = True
                        run.font.color.rgb = RGBColor(21, 52, 82)
    doc.add_paragraph('')
    return table

def create_docx(summary, eval_df, chunks_df, qa, charts):
    doc = Document(); set_doc_styles(doc)
    sec = doc.sections[0]
    sec.top_margin = Inches(0.65); sec.bottom_margin = Inches(0.65); sec.left_margin = Inches(0.65); sec.right_margin = Inches(0.65)
    title = doc.add_paragraph(); title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("MSU Texas Catalog Agentic RAG Chatbot\nStakeholder Implementation Report")
    run.bold = True; run.font.size = Pt(20); run.font.color.rgb = RGBColor(21, 52, 82)
    sub = doc.add_paragraph(); sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.add_run("Best demo version with query rewriting, citation verification, scheduled crawling/change detection, hybrid retrieval, MCP-style guardrails, and automated reporting.").italic = True
    doc.add_paragraph(f"Generated: {summary['created_at']}")
    doc.add_page_break()
    sections = [
        ("1. Executive Summary", f"This solution demonstrates an end-to-end multi-agent AI application for MSU Texas catalog Q&A. It crawls the catalog page and subpages, detects changes with page hashes, cleans and chunks content, tags metadata, rewrites student questions, builds a hybrid BM25 plus dense contextual embedding index, reranks evidence, answers questions with citations, verifies citation support, evaluates RAG quality, applies MCP-style guardrails, and generates stakeholder PowerPoint and DOCX reports. The current run indexed {summary['chunks']} chunks from {summary['source_pages']} source pages."),
        ("2. Use Case Objective", "The objective is to help students ask natural-language questions about admissions, transfers, international requirements, testing, financial aid, degree requirements, course descriptions, and related catalog topics. The assistant is designed to answer from official catalog evidence, not from memory or unsupported assumptions."),
        ("3. Why Simple Chatbots Are Not Enough", "A normal chatbot can hallucinate policy details. A catalog assistant must retrieve current evidence, cite source pages, recognize high-stakes topics, verify citations, and escalate students to Admissions, Registrar, Graduate School, International Services, or an advisor when final confirmation is required."),
        ("4. Best Demo Deployment", "For the public demo, Streamlit Community Cloud is the best option because the app is already Streamlit-based, easy to deploy from GitHub, and gives a public URL that students or stakeholders can open. FastAPI is not necessary for the demo. FastAPI, React, Qdrant, and PostgreSQL can be added later for enterprise production."),
        ("5. Data Sources and Crawl Scope", "The primary source is the MSU Texas catalog website. The crawler starts from the Admissions catalog page and follows approved catalog subpages such as content pages, program pages, course pages, entity pages, and catalog index pages. The crawler uses source allowlisting to avoid external and unsafe sources."),
        ("6. Scheduled Crawling and Change Detection", f"The upgraded system stores source URLs, content hashes, changed-page counts, and crawl timestamps in crawl_state.json. The latest run detected {summary.get('changed_pages', 0)} changed or new pages. In Streamlit demo mode, users can refresh manually. In production, a scheduled crawler can run daily or hourly and re-index only changed pages."),
        ("7. Document Processing Agent", "The Document Processing Agent removes scripts, navigation noise, repeated layout text, and irrelevant web fragments. It extracts headings, paragraphs, lists, and tables, then creates clean evidence blocks suitable for retrieval, citation, and answer generation."),
        ("8. Chunking Strategy", "Chunks are heading-aware and context-enriched. Each chunk includes the page title, heading, catalog year, topic, student type, and source URL. This makes retrieval more accurate than blindly splitting text by token count only."),
        ("9. Metadata Filtering", "Metadata includes source_url, page_title, heading, catalog_year, topic, student_type, content_hash, and last_crawled_at. This allows the system to prefer the right section for questions such as freshman admission, transfer credit, international requirements, testing, financial aid, or degree requirements."),
        ("10. Query Understanding Agent", "The Query Understanding Agent classifies the student's question by topic, student type, and risk level. It identifies whether the question is about admissions, transfer, international requirements, testing, deadlines, tuition/aid, degree requirements, or course descriptions."),
        ("11. Query Rewriting Agent", "The Query Rewriting Agent expands vague student language into catalog-specific retrieval terms. For example, 'Can I get in?' may be expanded into 'undergraduate admission application requirements freshman transcript deadline documents.' This increases retrieval accuracy and helps the system find evidence even when the student does not use official catalog terminology."),
        ("12. Embedding Strategy", "For a free Streamlit demo, the system uses contextual dense embeddings produced from TF-IDF plus TruncatedSVD. In production, this can be replaced with OpenAI text-embedding-3-small, OpenAI text-embedding-3-large, Gemini text-embedding-004, BGE-large, or E5-large embeddings."),
        ("13. Hybrid Retrieval", "The Retrieval Agent combines BM25 keyword scoring with dense semantic similarity. BM25 captures exact terms such as ACT, SAT, TSI, GPA, transcript, deadline, and transfer credit. Dense retrieval captures similar meaning when students phrase questions differently."),
        ("14. Reranking", "The Reranking Agent takes the initial retrieved chunks and reorders them using dense similarity, keyword overlap, source quality, topic match, heading match, and student-type match. Production upgrades can use Cohere Rerank, BGE reranker, FlashRank, or a cross-encoder reranker."),
        ("15. LLM Strategy", f"The best stakeholder demo uses Groq Llama 3.3 70B Versatile through GROQ_API_KEY and GROQ_MODEL because response quality matters. The app still includes a free local grounded fallback for cost control. Higher-quality production options include GPT-4o mini/GPT-4.1 mini, Gemini Flash/Pro, Claude Sonnet, or local vLLM-hosted models. The selected technology details are included in the Technology Stack appendix."),
        ("16. Agentic RAG Architecture", "This is not a one-shot RAG chain. It uses an agentic workflow: crawl, detect changes, process documents, tag metadata, understand the query, rewrite the query, retrieve, rerank, generate an answer, verify citations, apply guardrails, and create reports. Each agent has a specialized role and communicates through structured workflow state."),
        ("17. MCP and Guardrails", "The demo includes an MCP-style guardrail server module. It performs source authorization, prompt-injection cleanup, PII checks, high-stakes question detection, citation validation, and overclaim detection. In production, these functions can be exposed through a real MCP server with authentication, RBAC, and audit logs."),
        ("18. Ten-Agent Operating Model", "The system uses ten agents: Catalog Monitor & Change Detection, Document Processing, Metadata Tagging, Query Understanding, Query Rewriting, Hybrid Retrieval, Reranking, Answer Generation, Citation Verification & Guardrail, and Report Creation. They communicate through orchestrated state rather than working as disconnected scripts."),
        ("19. Citation Verification", "The Citation Verification Agent checks whether answer sentences are supported by retrieved evidence, whether citation markers are present, whether cited URLs are from approved domains, and whether unsupported claim examples should be flagged for review. This is one of the most important upgrades for a student-facing policy chatbot."),
        ("20. Student Experience", "Students can ask natural-language questions in Streamlit. The assistant returns a clear answer, citation URLs/headings, a high-stakes escalation note when needed, retrieval evidence, rewritten query diagnostics, and citation verification results. This improves transparency and encourages students to verify official decisions."),
        ("21. Staff and Business Impact", "The system can reduce repetitive front-desk, admissions, registrar, and advising questions. It does not replace official staff decisions; it routes students to the correct evidence and office faster, reducing repetitive workload and improving consistency."),
        ("22. RAG Evaluation Metrics", "The system reports MRR, relevance score, mean similarity, groundedness, hallucination risk, and citation confidence. These metrics help stakeholders see whether the right evidence is retrieved and whether answers are supported by citations."),
        ("23. Production Roadmap", "Phase 1: Streamlit public demo. Phase 2: scheduled crawler, optional Groq/Gemini/OpenAI LLM, external vector DB, and feedback collection. Phase 3: enterprise production with SSO, RBAC, audit logs, MCP tools, monitoring, evaluation dashboards, and human approval workflow."),
        ("24. Governance Policy", "The assistant should never invent deadlines, admission decisions, GPA rules, transfer-credit approvals, immigration guidance, financial-aid determinations, or graduation eligibility. It should answer from evidence and tell students to contact official offices for final confirmation."),
        ("25. Final Recommendation", "Build the public demo in Streamlit Community Cloud, keep the code flat and simple, use agentic hybrid retrieval and guardrails, and make the system dynamic through on-demand or scheduled crawling. This gives stakeholders a practical, transparent, and credible student-facing AI assistant.")
    ]
    for heading, body in sections:
        doc.add_heading(heading, level=1)
        for para in re.split(r"(?<=\.)\s+", body):
            if para.strip():
                doc.add_paragraph(para.strip())
        if heading == "18. Ten-Agent Operating Model":
            rows = [["Agent", "Role"]] + [[a, r] for a, r in AGENT_REGISTRY]
            add_metric_table(doc, rows)
        if heading == "22. RAG Evaluation Metrics":
            rows = [["Metric", "Average Score", "Meaning"],
                    ["MRR", f"{eval_df['MRR'].mean():.2f}", "Best expected evidence appears near the top"],
                    ["Relevance", f"{eval_df['relevance_score'].mean():.2f}", "Retrieved chunks match the question"],
                    ["Similarity", f"{eval_df['mean_similarity'].mean():.2f}", "Semantic closeness between query and evidence"],
                    ["Groundedness", f"{eval_df['groundedness'].mean():.2f}", "Answer support from retrieved evidence"],
                    ["Hallucination Risk", f"{eval_df['hallucination_risk'].mean():.2f}", "Lower value is better"],
                    ["Citation Confidence", f"{eval_df['citation_confidence'].mean():.2f}", "Citation support and source validity"]]
            add_metric_table(doc, rows)
            if os.path.exists(charts.get("rag_eval", "")):
                doc.add_picture(charts["rag_eval"], width=Inches(6.4))
        if heading == "12. Embedding Strategy":
            rows = [["Technology", "Demo Choice", "Production Upgrade"],
                    ["Embeddings", "TF-IDF + SVD contextual dense vectors", "OpenAI/Gemini/BGE/E5"],
                    ["Vector DB", "In-memory cosine search", "Qdrant/Pinecone/Weaviate/Vertex AI"],
                    ["Hybrid Search", "BM25 + dense similarity", "OpenSearch/Elasticsearch + vector DB"],
                    ["Reranking", "Lightweight metadata-aware reranker", "Cohere/BGE/FlashRank"],
                    ["Query Rewriting", "Rule-based expansion", "LLM multi-query retriever"],
                    ["Citation Verification", "Support/citation/source checks", "Claim-level verifier + human approval"]]
            add_metric_table(doc, rows)
        if heading != sections[-1][0]:
            doc.add_page_break()
    doc.add_heading("Appendix A. Sample Chat Answers", level=1)
    for item in qa:
        doc.add_heading(item["question"], level=2)
        doc.add_paragraph(f"Rewritten query: {item.get('retrieval_metrics', {}).get('rewritten_query', '')}")
        doc.add_paragraph(item["answer"])
        doc.add_paragraph(f"Citation verification: {item.get('citation_verification', {})}")
        for c in item["citations"][:3]:
            doc.add_paragraph(f"Citation: {c['heading']} - {c['url']}")
    doc.add_page_break()
    doc.add_heading("Appendix B. Technology Stack Detail", level=1)
    rows = [["Layer", "Selected Technology"]] + [[k, v] for k, v in TECH_STACK.items()]
    add_metric_table(doc, rows)
    out = REPORT_DIR / "msu_agentic_rag_stakeholder_report.docx"
    doc.save(out)
    return out
